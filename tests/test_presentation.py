"""Тесты просмотрщика презентаций: PPTX / PPT / ODP.

Поведение проверяется через главный шов ``ViewerRegistry`` и публичный API
``PresentationViewer``:

* реестр отдаёт ``PresentationViewer`` по ``pptx`` / ``ppt`` / ``odp``;
* PPTX — текст слайдов в порядке следования, разделённый по слайдам;
* вложенные в слайды изображения доходят до виджета;
* legacy ``.ppt`` открывается хотя бы как извлечённый текст;
* ODP открывается как текст слайдов;
* битый образец → «ошибочный» виджет, без падения.
"""

from __future__ import annotations

import io
import sys
import zlib
from pathlib import Path

import pytest
from PyQt6.QtGui import QTextDocument
from PyQt6.QtWidgets import QApplication

from omniviewer.registry import ViewerRegistry
from omniviewer.viewers.presentation import PresentationViewer

DEMO = Path(__file__).parent.parent / "demo"
_IMG_RESOURCE = QTextDocument.ResourceType.ImageResource.value


@pytest.fixture(scope="session", autouse=True)
def qapp():
    app = QApplication.instance() or QApplication(sys.argv)
    yield app


@pytest.fixture
def registry() -> ViewerRegistry:
    return ViewerRegistry()


# --------------------------------------------------------------------------- #
# Помощники — строим образцы прямо в тесте, без зависимости от demo/           #
# --------------------------------------------------------------------------- #


def _png(w: int = 12, h: int = 12) -> bytes:
    import struct

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    rows = bytearray()
    for _y in range(h):
        rows.append(0)
        rows += bytes((200, 80, 40)) * w
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", zlib.compress(bytes(rows), 9))
        + chunk(b"IEND", b"")
    )


def _make_pptx(*, with_picture: bool = True) -> bytes:
    from pptx import Presentation
    from pptx.util import Emu, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    tb1 = s1.shapes.add_textbox(Emu(457200), Emu(457200), Emu(8000000), Emu(900000))
    tf = tb1.text_frame
    tf.text = "ПЕРВЫЙ СЛАЙД"
    tf.add_paragraph().text = "строка один"
    tf.add_paragraph().text = "строка два"

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Emu(457200), Emu(457200), Emu(8000000), Emu(900000))
    tb2.text_frame.text = "ВТОРОЙ СЛАЙД"
    if with_picture:
        s2.shapes.add_picture(io.BytesIO(_png()), Emu(457200), Emu(2000000), height=Pt(60))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_odp() -> bytes:
    from odf.draw import Frame, Page, TextBox
    from odf.opendocument import OpenDocumentPresentation
    from odf.style import MasterPage, PageLayout, PageLayoutProperties
    from odf.text import P

    doc = OpenDocumentPresentation()
    layout = PageLayout(name="pl1")
    doc.automaticstyles.addElement(layout)
    layout.addElement(PageLayoutProperties(pagewidth="28cm", pageheight="21cm"))
    master = MasterPage(name="Default", pagelayoutname=layout)
    doc.masterstyles.addElement(master)

    for name, lines in (
        ("Слайд один", ["ЗАГОЛОВОК ODP", "вводная строка"]),
        ("Слайд два", ["ВТОРОЙ ЛИСТ", "пункт A", "пункт B"]),
    ):
        page = Page(name=name, masterpagename=master)
        doc.presentation.addElement(page)
        frame = Frame(width="24cm", height="12cm", x="1cm", y="1cm")
        page.addElement(frame)
        tb = TextBox()
        frame.addElement(tb)
        for line in lines:
            tb.addElement(P(text=line))

    buf = io.BytesIO()
    doc.write(buf)
    return buf.getvalue()


# --------------------------------------------------------------------------- #
# Диспетчеризация                                                             #
# --------------------------------------------------------------------------- #


def test_registry_selects_presentation_viewer_pptx(registry, tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    p.write_bytes(_make_pptx())
    assert isinstance(registry.viewer_for(p), PresentationViewer)


def test_registry_selects_presentation_viewer_odp(registry, tmp_path: Path) -> None:
    p = tmp_path / "deck.odp"
    p.write_bytes(_make_odp())
    assert isinstance(registry.viewer_for(p), PresentationViewer)


def test_registry_selects_presentation_viewer_ppt(registry) -> None:
    sample = DEMO / "presentations/sample.ppt"
    if not sample.exists():
        pytest.skip("demo/presentations/sample.ppt не сгенерирован")
    assert isinstance(registry.viewer_for(sample), PresentationViewer)


def test_presentation_viewer_ignores_plain_zip_and_text(registry, tmp_path: Path) -> None:
    import zipfile

    z = tmp_path / "plain.zip"
    with zipfile.ZipFile(z, "w") as zf:
        zf.writestr("a.txt", b"hi")
    assert not isinstance(registry.viewer_for(z), PresentationViewer)

    t = tmp_path / "notes.txt"
    t.write_bytes(b"just text")
    assert not isinstance(registry.viewer_for(t), PresentationViewer)


# --------------------------------------------------------------------------- #
# PPTX                                                                        #
# --------------------------------------------------------------------------- #


def test_pptx_slide_text_in_order_split_by_slides(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    p.write_bytes(_make_pptx())
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert not viewer.is_error_widget
    assert viewer.slide_count == 2

    plain = viewer._browser.toPlainText()
    assert "ПЕРВЫЙ СЛАЙД" in plain
    assert "строка один" in plain and "строка два" in plain
    assert "ВТОРОЙ СЛАЙД" in plain
    # порядок слайдов сохранён
    assert plain.index("ПЕРВЫЙ СЛАЙД") < plain.index("ВТОРОЙ СЛАЙД")
    # текст первого слайда не перемешан со вторым
    assert plain.index("строка два") < plain.index("ВТОРОЙ СЛАЙД")


def test_pptx_embedded_image_reaches_widget(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    p.write_bytes(_make_pptx(with_picture=True))
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert not viewer.is_error_widget
    assert viewer.image_count >= 1
    # каждый зарегистрированный ресурс отдаётся браузером как непустой блоб
    for key in viewer.resource_keys:
        from PyQt6.QtCore import QUrl

        data = viewer._browser.loadResource(_IMG_RESOURCE, QUrl(key))
        assert data is not None and len(bytes(data)) > 0


def test_pptx_without_images_has_no_resources(tmp_path: Path) -> None:
    p = tmp_path / "deck.pptx"
    p.write_bytes(_make_pptx(with_picture=False))
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert not viewer.is_error_widget
    assert viewer.image_count == 0


def test_truncated_pptx_yields_error_widget(tmp_path: Path) -> None:
    full = _make_pptx()
    p = tmp_path / "broken.pptx"
    p.write_bytes(full[: len(full) // 3])
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert viewer.is_error_widget


# --------------------------------------------------------------------------- #
# ODP                                                                         #
# --------------------------------------------------------------------------- #


def test_odp_opens_as_slide_text(tmp_path: Path) -> None:
    p = tmp_path / "deck.odp"
    p.write_bytes(_make_odp())
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert not viewer.is_error_widget
    assert viewer.slide_count == 2
    plain = viewer._browser.toPlainText()
    assert "ЗАГОЛОВОК ODP" in plain
    assert "пункт A" in plain and "пункт B" in plain
    assert plain.index("ЗАГОЛОВОК ODP") < plain.index("ВТОРОЙ ЛИСТ")


# --------------------------------------------------------------------------- #
# Legacy PPT                                                                  #
# --------------------------------------------------------------------------- #


def test_ppt_legacy_opens_as_extracted_text() -> None:
    sample = DEMO / "presentations/sample.ppt"
    if not sample.exists():
        pytest.skip("demo/presentations/sample.ppt не сгенерирован")
    viewer = PresentationViewer()
    viewer.safe_load(sample)

    assert not viewer.is_error_widget
    plain = viewer._browser.toPlainText()
    # текстовые атомы из потока «PowerPoint Document» (Unicode и ANSI)
    assert "Слайд 1" in plain
    assert "Slide 2 body text" in plain


def test_ppt_garbage_ole_yields_error_widget(tmp_path: Path) -> None:
    # корректная OLE-сигнатура, но никакого потока PowerPoint Document
    p = tmp_path / "weird.ppt"
    p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 4088)
    viewer = PresentationViewer()
    viewer.safe_load(p)

    assert viewer.is_error_widget
