"""Просмотрщик презентаций: PPTX / PPT / ODP — «показать содержание».

Единый конвейер «презентация → строка HTML → QTextBrowser» (см.
:mod:`omniviewer.viewers.html_render`). Каждый слайд — секция HTML: подзаголовок
«Слайд N», далее текст слайда в порядке следования фигур и встроенные в него
изображения. Legacy ``.ppt`` (двоичный формат) открывается хотя бы как
извлечённый текст. Рендер строго офлайн: внешние ресурсы не загружаются.

Форматы и движки:

* ``pptx`` — :mod:`python-pptx` (текст фигур по порядку + блобы картинок);
* ``odp``  — :mod:`odfpy` (текст страниц-слайдов, картинки из ``Pictures/``);
* ``ppt``  — :mod:`olefile` + разбор записей потока «PowerPoint Document»
  (``TextCharsAtom`` / ``TextBytesAtom``).
"""

from __future__ import annotations

import html as _html
import struct
from pathlib import Path

from PyQt6.QtCore import QMimeType

from omniviewer.viewers.base import BaseViewer
from omniviewer.viewers.html_render import build_html_browser

_PPTX_SUFFIXES = (".pptx", ".pptm", ".ppsx", ".ppsm")
_PPT_SUFFIXES = (".ppt", ".pps")
_ODP_SUFFIXES = (".odp", ".otp")

_PPTX_MIMES = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
        "application/vnd.ms-powerpoint.presentation.macroenabled.12",
        "application/vnd.ms-powerpoint.slideshow.macroenabled.12",
    }
)
_PPT_MIMES = frozenset(
    {
        "application/vnd.ms-powerpoint",
        "application/mspowerpoint",
        "application/x-mspowerpoint",
    }
)
_ODP_MIMES = frozenset(
    {
        "application/vnd.oasis.opendocument.presentation",
        "application/vnd.oasis.opendocument.presentation-template",
    }
)

_ODP_MIMETYPE = b"application/vnd.oasis.opendocument.presentation"
_OLE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

# Типы записей двоичного PowerPoint (MS-PPT).
_RT_TEXT_CHARS_ATOM = 0x0FA0  # UTF-16LE
_RT_TEXT_BYTES_ATOM = 0x0FA8  # ANSI (одна кодовая страница, берём latin-1)
_RT_CSTRING = 0x0FBA  # UTF-16LE, встречается в заголовках/заметках


def _ppt_text_atoms(stream: bytes) -> list[str]:
    """Собрать текст из записей потока «PowerPoint Document».

    Формат записи: 8-байтовый заголовок ``recVer/recInstance`` (uint16),
    ``recType`` (uint16), ``recLen`` (uint32), затем тело. Контейнеры помечены
    ``recVer == 0xF`` — в них рекурсируемся; текстовые атомы декодируем.
    """
    out: list[str] = []

    def walk(buf: bytes, depth: int) -> None:
        i, n = 0, len(buf)
        while i + 8 <= n:
            ver_inst, rec_type, rec_len = struct.unpack_from("<HHI", buf, i)
            i += 8
            if rec_len > n - i:
                break
            body = buf[i : i + rec_len]
            i += rec_len
            if rec_type in (_RT_TEXT_CHARS_ATOM, _RT_CSTRING):
                out.append(body.decode("utf-16-le", "replace"))
            elif rec_type == _RT_TEXT_BYTES_ATOM:
                out.append(body.decode("latin-1", "replace"))
            elif (ver_inst & 0x0F) == 0x0F and depth < 24:
                walk(body, depth + 1)

    walk(stream, 0)
    return out


## @brief Просмотрщик презентаций (PPTX / PPT / ODP).
#
# Приводит презентацию к строке HTML (каждый слайд — секция «Слайд N» с его
# текстом и картинками) и показывает во встроенном движке rich text Qt
# (QTextBrowser) через общий офлайн-хелпер. Legacy ``.ppt`` показывается как
# извлечённый из двоичного потока текст. Внешние сетевые ресурсы не грузятся.
class PresentationViewer(BaseViewer):
    mime_types = tuple(_PPTX_MIMES | _PPT_MIMES | _ODP_MIMES)
    extensions = _PPTX_SUFFIXES + _PPT_SUFFIXES + _ODP_SUFFIXES
    priority = 30  # выше ArchiveViewer(15): pptx/odp — zip-контейнеры

    def __init__(self):
        super().__init__()
        self._browser = build_html_browser("")
        self._layout.addWidget(self._browser)
        self.rendered_html: str = ""
        self.slide_count: int = 0
        self.image_count: int = 0
        self.resource_keys: list[str] = []

    @classmethod
    def can_handle(cls, path: Path, mime: QMimeType | str) -> bool:
        if path.suffix.lower() in cls.extensions:
            return True
        mime_name = mime.name() if isinstance(mime, QMimeType) else str(mime)
        return mime_name in _PPTX_MIMES or mime_name in _PPT_MIMES or mime_name in _ODP_MIMES

    def load(self, path: Path) -> None:
        kind = self._detect_kind(path)
        if kind == "odp":
            slides, resources = self._load_odp(path)
            self._render(slides, resources)
        elif kind == "ppt":
            frags = self._load_ppt(path)
            self._render([frags], {}, numbered=False, lead="Извлечённый текст (legacy .ppt)")
        else:
            slides, resources = self._load_pptx(path)
            self._render(slides, resources)

    # ------------------------------------------------------------------ #
    # Определение формата                                                #
    # ------------------------------------------------------------------ #

    @staticmethod
    def _detect_kind(path: Path) -> str:
        suffix = path.suffix.lower()
        if suffix in _ODP_SUFFIXES:
            return "odp"
        if suffix in _PPT_SUFFIXES:
            return "ppt"
        if suffix in _PPTX_SUFFIXES:
            return "pptx"
        head = path.read_bytes()[:8]
        if head.startswith(_OLE_MAGIC):
            return "ppt"
        if head.startswith(b"PK\x03\x04"):
            import zipfile

            try:
                with zipfile.ZipFile(path) as zf:
                    if zf.read("mimetype").strip() == _ODP_MIMETYPE:
                        return "odp"
            except (KeyError, OSError, zipfile.BadZipFile):
                pass
            return "pptx"
        raise ValueError("Не удалось определить формат презентации")

    # ------------------------------------------------------------------ #
    # Загрузчики форматов                                                #
    # ------------------------------------------------------------------ #

    def _load_pptx(self, path: Path) -> tuple[list[list[str]], dict[str, bytes]]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(path))
        slides: list[list[str]] = []
        resources: dict[str, bytes] = {}

        def emit_shape(shape, frags: list[str]) -> None:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    emit_shape(child, frags)
                return
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = ("".join(run.text for run in para.runs) or para.text).strip()
                    if text:
                        frags.append(f"<p>{_html.escape(text)}</p>")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ext = (shape.image.ext or "png").lstrip(".")
                except Exception:  # noqa: BLE001 - картинка без извлекаемого блоба
                    return
                key = f"pptx-img-{len(resources) + 1}.{ext}"
                resources[key] = blob
                frags.append(f'<p><img src="{key}"></p>')

        for slide in prs.slides:
            frags: list[str] = []
            for shape in slide.shapes:
                emit_shape(shape, frags)
            slides.append(frags)
        return slides, resources

    def _load_odp(self, path: Path) -> tuple[list[list[str]], dict[str, bytes]]:
        from odf import teletype
        from odf.draw import Image as DrawImage
        from odf.draw import Page
        from odf.namespaces import TEXTNS
        from odf.opendocument import load

        doc = load(str(path))
        pictures = getattr(doc, "Pictures", {}) or {}
        para_qnames = {(TEXTNS, "p"), (TEXTNS, "h")}

        def collect(node, acc: list[str]) -> None:
            """Собрать текст абзацев/заголовков поддерева в порядке обхода."""
            for child in getattr(node, "childNodes", []):
                if getattr(child, "qname", None) in para_qnames:
                    text = teletype.extractText(child).strip()
                    if text:
                        acc.append(f"<p>{_html.escape(text)}</p>")
                else:
                    collect(child, acc)

        slides: list[list[str]] = []
        resources: dict[str, bytes] = {}

        for page in doc.getElementsByType(Page):
            frags: list[str] = []
            collect(page, frags)

            for img in page.getElementsByType(DrawImage):
                href = img.getAttribute("href")
                entry = pictures.get(href) if href else None
                if not entry:
                    continue
                content = entry[1] if isinstance(entry, tuple) else entry
                if not content:
                    continue
                key = f"odp-img-{len(resources) + 1}"
                resources[key] = bytes(content)
                frags.append(f'<p><img src="{key}"></p>')

            slides.append(frags)
        return slides, resources

    def _load_ppt(self, path: Path) -> list[str]:
        import olefile

        if not olefile.isOleFile(str(path)):
            raise ValueError("Файл .ppt не является OLE-документом")
        ole = olefile.OleFileIO(str(path))
        try:
            if not ole.exists("PowerPoint Document"):
                raise ValueError("В .ppt нет потока «PowerPoint Document»")
            raw = ole.openstream("PowerPoint Document").read()
        finally:
            ole.close()

        chunks = _ppt_text_atoms(raw)
        frags: list[str] = []
        for chunk in chunks:
            normalized = chunk.replace("\r", "\n").replace("\x0b", "\n").replace("\x00", "")
            for line in normalized.split("\n"):
                line = line.strip()
                if line:
                    frags.append(f"<p>{_html.escape(line)}</p>")
        if not frags:
            raise ValueError("Не удалось извлечь текст из legacy .ppt")
        return frags

    # ------------------------------------------------------------------ #
    # Рендер                                                             #
    # ------------------------------------------------------------------ #

    def _render(
        self,
        slides: list[list[str]],
        resources: dict[str, bytes],
        *,
        numbered: bool = True,
        lead: str | None = None,
    ) -> None:
        parts = ['<meta charset="utf-8">']
        if lead:
            parts.append(f"<p><i>{_html.escape(lead)}</i></p>")
        for i, frags in enumerate(slides, 1):
            if numbered:
                parts.append(f"<h2>Слайд {i}</h2>")
            parts.append("".join(frags) if frags else "<p><i>(без текста)</i></p>")
            if i != len(slides):
                parts.append("<hr>")
        html = "\n".join(parts)

        self.rendered_html = html
        self.slide_count = len(slides)
        self.image_count = len(resources)
        self.resource_keys = list(resources)
        self._browser.set_resources(resources)
        self._browser.setHtml(html)
